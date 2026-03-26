"""
Generate: Masterclass Interpretability (EN) — clean rebuild
Run with: uv run --with python-pptx generate_slides.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Brand colors ───────────────────────────────────────────────────────────────
GREEN       = RGBColor(0x75, 0xDF, 0xC1)   # DataScientest teal-green
NAVY        = RGBColor(0x1A, 0x2E, 0x4A)   # dark navy background
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY  = RGBColor(0xF4, 0xF5, 0xF7)
MID_GRAY    = RGBColor(0x6B, 0x7A, 0x8D)
ORANGE      = RGBColor(0xFF, 0x6B, 0x35)   # warnings / limitations
YELLOW      = RGBColor(0xFF, 0xD1, 0x66)   # highlights
DARK_GRAY   = RGBColor(0x2D, 0x3A, 0x4A)

# ── Slide dimensions (16:9 widescreen) ────────────────────────────────────────
W = Inches(13.33)
H = Inches(7.5)


# ══════════════════════════════════════════════════════════════════════════════
# Helpers
# ══════════════════════════════════════════════════════════════════════════════

def new_prs():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs


def blank_slide(prs):
    blank_layout = prs.slide_layouts[6]   # truly blank
    return prs.slides.add_slide(blank_layout)


def fill_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, x, y, w, h, fill_color, border_color=None, border_pt=0):
    shape = slide.shapes.add_shape(1, x, y, w, h)   # MSO_SHAPE_TYPE.RECTANGLE = 1
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color and border_pt > 0:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(border_pt)
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, text, x, y, w, h,
             font_size=18, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, wrap=True, italic=False):
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox


def add_text_box(slide, text, x, y, w, h,
                 font_size=18, bold=False, color=WHITE,
                 align=PP_ALIGN.LEFT, bg_color=None, border_color=None,
                 italic=False, padding_pt=6):
    """Text inside a filled rectangle (card style)."""
    if bg_color:
        add_rect(slide, x, y, w, h, bg_color, border_color, 1 if border_color else 0)
    txBox = slide.shapes.add_textbox(
        x + Pt(padding_pt), y + Pt(padding_pt),
        w - Pt(padding_pt*2), h - Pt(padding_pt*2)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox


def add_multiline(slide, lines, x, y, w, h,
                  font_size=16, color=WHITE, bold_first=False,
                  line_spacing_pt=None, align=PP_ALIGN.LEFT):
    """Add a text box with multiple paragraphs."""
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.alignment = align
        if line_spacing_pt:
            from pptx.util import Pt as _Pt
            p.space_before = _Pt(line_spacing_pt)
        run = p.add_run()
        run.text = line
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
        run.font.bold = (bold_first and i == 0)
    return txBox


def section_divider(prs, title, subtitle=""):
    """Navy slide with centered section title."""
    slide = blank_slide(prs)
    fill_bg(slide, NAVY)
    # Green accent bar left
    add_rect(slide, Inches(0), Inches(2.8), Inches(0.18), Inches(1.9), GREEN)
    add_text(slide, title,
             Inches(0.5), Inches(3.0), Inches(12), Inches(1.2),
             font_size=40, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    if subtitle:
        add_text(slide, subtitle,
                 Inches(0.5), Inches(4.3), Inches(12), Inches(0.7),
                 font_size=20, color=GREEN, align=PP_ALIGN.CENTER)
    return slide


def content_slide(prs, title):
    """White slide with navy title bar."""
    slide = blank_slide(prs)
    fill_bg(slide, WHITE)
    # Title bar
    add_rect(slide, Inches(0), Inches(0), W, Inches(1.1), NAVY)
    add_text(slide, title,
             Inches(0.35), Inches(0.12), Inches(12.6), Inches(0.85),
             font_size=26, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    # Green bottom stripe
    add_rect(slide, Inches(0), H - Inches(0.12), W, Inches(0.12), GREEN)
    return slide


def bullet_card(slide, items, x, y, w, h, bg=LIGHT_GRAY, title=None, title_color=NAVY):
    """A card with optional title and bullet list."""
    add_rect(slide, x, y, w, h, bg)
    offset = Inches(0.1)
    if title:
        add_text(slide, title,
                 x + offset, y + Inches(0.1), w - offset*2, Inches(0.4),
                 font_size=16, bold=True, color=title_color)
        y += Inches(0.5)
        h -= Inches(0.5)
    add_multiline(slide, items, x + offset, y + Inches(0.05),
                  w - offset*2, h - Inches(0.1),
                  font_size=14, color=DARK_GRAY)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDES
# ══════════════════════════════════════════════════════════════════════════════

def slide_title(prs):
    """Slide 1 — Title"""
    slide = blank_slide(prs)
    fill_bg(slide, NAVY)
    # Big green accent block left
    add_rect(slide, Inches(0), Inches(0), Inches(0.55), H, GREEN)
    # Subtle horizontal line
    add_rect(slide, Inches(0.7), Inches(3.5), Inches(12.4), Inches(0.04), GREEN)
    add_text(slide, "MASTERCLASS",
             Inches(0.9), Inches(1.5), Inches(11), Inches(1.0),
             font_size=20, bold=False, color=GREEN, align=PP_ALIGN.LEFT)
    add_text(slide, "Model Interpretability",
             Inches(0.9), Inches(2.3), Inches(11), Inches(1.4),
             font_size=44, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    add_text(slide, "Understanding why your model predicts what it predicts",
             Inches(0.9), Inches(3.7), Inches(10), Inches(0.7),
             font_size=18, color=MID_GRAY, align=PP_ALIGN.LEFT)
    add_text(slide, "DataScientest",
             Inches(0.9), Inches(6.5), Inches(5), Inches(0.5),
             font_size=14, color=MID_GRAY)


def slide_housekeeping(prs):
    """Slide 2 — Housekeeping"""
    slide = content_slide(prs, "Before We Start")
    rules = [
        ("📸", "Validate your presence",       "Make sure you're marked present in the platform"),
        ("📷", "Activate your camera",          "Helps engagement and interaction — thank you!"),
        ("✋", "Respect communication rules",   "Use the raise-hand feature to ask questions"),
        ("🔇", "Mute when not speaking",        "Reduces background noise for everyone"),
    ]
    for i, (icon, rule, detail) in enumerate(rules):
        col = i % 2
        row = i // 2
        x = Inches(0.4) + col * Inches(6.4)
        y = Inches(1.4) + row * Inches(2.5)
        add_rect(slide, x, y, Inches(6.1), Inches(2.2), LIGHT_GRAY)
        add_text(slide, icon,
                 x + Inches(0.15), y + Inches(0.1), Inches(0.8), Inches(0.8),
                 font_size=30, align=PP_ALIGN.CENTER)
        add_text(slide, rule,
                 x + Inches(1.0), y + Inches(0.15), Inches(4.8), Inches(0.5),
                 font_size=16, bold=True, color=NAVY)
        add_text(slide, detail,
                 x + Inches(1.0), y + Inches(0.65), Inches(4.8), Inches(1.3),
                 font_size=13, color=MID_GRAY)


def slide_agenda(prs):
    """Slide 3 — Agenda"""
    slide = content_slide(prs, "Agenda")
    sections = [
        ("01", "The Problem",        "A real-world case study — why does interpretability matter?"),
        ("02", "What & Why",         "Definitions, dimensions, and motivation"),
        ("03", "Traditional Methods","Feature importance, coefficients, PCA — and their limits"),
        ("04", "SHAP",               "The modern standard: global + local explanations"),
        ("05", "Hands-on Notebook",  "Apply SHAP on census data and text classification"),
        ("06", "Your Turn",          "Open challenge — interpret a model of your choice"),
    ]
    for i, (num, title, desc) in enumerate(sections):
        col = i % 2
        row = i // 2
        x = Inches(0.35) + col * Inches(6.45)
        y = Inches(1.25) + row * Inches(1.95)
        add_rect(slide, x, y, Inches(6.15), Inches(1.75), LIGHT_GRAY)
        add_rect(slide, x, y, Inches(0.6), Inches(1.75), GREEN)
        add_text(slide, num,
                 x + Inches(0.02), y + Inches(0.55), Inches(0.56), Inches(0.65),
                 font_size=18, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        add_text(slide, title,
                 x + Inches(0.75), y + Inches(0.15), Inches(5.2), Inches(0.5),
                 font_size=16, bold=True, color=NAVY)
        add_text(slide, desc,
                 x + Inches(0.75), y + Inches(0.65), Inches(5.2), Inches(0.95),
                 font_size=13, color=MID_GRAY)


def slide_estelle_1(prs):
    """Slide 4 — Estelle case study setup"""
    slide = content_slide(prs, "A Real Story: Estelle's Loan Application")
    # Left: story
    add_multiline(slide, [
        "Estelle is looking for her first real estate investment.",
        "",
        "She goes to her bank to apply for a loan.",
        "",
        "She provides all the required documents:",
        "  • Income history",
        "  • Employment status",
        "  • Credit score",
        "  • Personal financial data",
    ], Inches(0.4), Inches(1.25), Inches(7.2), Inches(5.8),
       font_size=17, color=DARK_GRAY)
    # Right: visual card
    add_rect(slide, Inches(8.0), Inches(1.3), Inches(5.0), Inches(5.5), LIGHT_GRAY)
    add_text(slide, "🏦",
             Inches(9.8), Inches(1.8), Inches(1.4), Inches(1.2),
             font_size=52, align=PP_ALIGN.CENTER)
    add_text(slide, "Bank Loan Application",
             Inches(8.1), Inches(3.1), Inches(4.8), Inches(0.5),
             font_size=16, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    fields = ["Age: 29", "Income: €32,000/yr", "Marital status: Single",
              "Employment: 4 years", "Capital savings: €2,400"]
    for i, f in enumerate(fields):
        y = Inches(3.7) + i * Inches(0.44)
        add_rect(slide, Inches(8.3), y, Inches(4.5), Inches(0.38),
                 RGBColor(0xE0, 0xF7, 0xF2))
        add_text(slide, f, Inches(8.45), y + Inches(0.04), Inches(4.2), Inches(0.34),
                 font_size=13, color=DARK_GRAY)


def slide_estelle_2(prs):
    """Slide 5 — Estelle: the rejection"""
    slide = content_slide(prs, "The Problem: High Performance, Zero Explanation")
    # Left panel
    add_rect(slide, Inches(0.35), Inches(1.25), Inches(6.0), Inches(5.5), LIGHT_GRAY)
    add_text(slide, "What the model does:",
             Inches(0.55), Inches(1.4), Inches(5.6), Inches(0.45),
             font_size=16, bold=True, color=NAVY)
    add_multiline(slide, [
        "✅  Trained on 50,000 loan cases",
        "✅  Accuracy: 91%",
        "✅  Processes Estelle's file in 0.3 seconds",
        "❌  REJECTED — ineligible for loan",
    ], Inches(0.55), Inches(1.95), Inches(5.6), Inches(2.8),
       font_size=15, color=DARK_GRAY)
    # Divider
    add_rect(slide, Inches(0.55), Inches(4.5), Inches(5.6), Inches(0.03), GREEN)
    add_text(slide, "Estelle asks: \"Why was I rejected?\"",
             Inches(0.55), Inches(4.65), Inches(5.6), Inches(0.55),
             font_size=15, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    add_text(slide, "Banker: \"The algorithm decided so.\nI cannot explain it.\"",
             Inches(0.55), Inches(5.25), Inches(5.6), Inches(0.9),
             font_size=15, color=ORANGE, align=PP_ALIGN.CENTER)

    # Right panel
    add_rect(slide, Inches(6.8), Inches(1.25), Inches(6.15), Inches(5.5), NAVY)
    add_text(slide, "Why this is a problem",
             Inches(7.0), Inches(1.45), Inches(5.7), Inches(0.5),
             font_size=17, bold=True, color=GREEN)
    issues = [
        ("⚖️  Legal", "GDPR Article 22: right to explanation\nfor automated decisions"),
        ("🤝  Trust", "Customers lose trust in decisions\nthey cannot understand"),
        ("🐛  Bugs", "Model may have learned spurious\ncorrelations — without interpretability,\nyou'll never know"),
        ("📊  Business", "Stakeholders need confidence\nbefore deploying a model"),
    ]
    for i, (icon, detail) in enumerate(issues):
        y = Inches(2.05) + i * Inches(1.1)
        add_text(slide, icon,
                 Inches(7.0), y, Inches(1.0), Inches(0.5),
                 font_size=15, color=GREEN)
        add_text(slide, detail,
                 Inches(8.0), y, Inches(4.7), Inches(0.95),
                 font_size=13, color=WHITE)


def slide_dimensions(prs):
    """Slide 6 — Two dimensions of interpretability"""
    slide = content_slide(prs, "Two Key Dimensions of Interpretability")

    # Scope axis
    add_rect(slide, Inches(0.4), Inches(1.3), Inches(6.1), Inches(2.7), LIGHT_GRAY)
    add_rect(slide, Inches(0.4), Inches(1.3), Inches(6.1), Inches(0.5), NAVY)
    add_text(slide, "SCOPE — What level does the explanation target?",
             Inches(0.55), Inches(1.33), Inches(5.8), Inches(0.44),
             font_size=14, bold=True, color=GREEN)
    add_multiline(slide, [
        "🌍  GLOBAL — How does the model behave overall?",
        "     Which features matter most across ALL predictions?",
        "     → Beeswarm plot, feature importance bar",
        "",
        "🔍  LOCAL — Why this specific prediction?",
        "     For this one person, which features drove the decision?",
        "     → Force plot, waterfall plot",
    ], Inches(0.55), Inches(1.9), Inches(5.8), Inches(2.0),
       font_size=13, color=DARK_GRAY)

    # Method type axis
    add_rect(slide, Inches(6.85), Inches(1.3), Inches(6.1), Inches(2.7), LIGHT_GRAY)
    add_rect(slide, Inches(6.85), Inches(1.3), Inches(6.1), Inches(0.5), NAVY)
    add_text(slide, "METHOD TYPE — When is the explanation generated?",
             Inches(7.0), Inches(1.33), Inches(5.8), Inches(0.44),
             font_size=14, bold=True, color=GREEN)
    add_multiline(slide, [
        "🔓  INTRINSIC — Built into the model",
        "     Simple, transparent structure",
        "     → Decision trees, linear regression",
        "",
        "🔬  POST HOC — Applied after training",
        "     Explains any black-box model",
        "     → SHAP, LIME",
    ], Inches(7.0), Inches(1.9), Inches(5.8), Inches(2.0),
       font_size=13, color=DARK_GRAY)

    # Journey arrow
    add_rect(slide, Inches(0.4), Inches(4.3), Inches(12.55), Inches(0.04), GREEN)
    add_text(slide, "Our journey in this masterclass  →",
             Inches(0.4), Inches(4.4), Inches(6.0), Inches(0.4),
             font_size=13, color=MID_GRAY)

    journey = [
        ("Decision Tree", "Intrinsic\nGlobal + Local", LIGHT_GRAY, NAVY),
        ("Log. Regression", "Intrinsic\nGlobal", LIGHT_GRAY, NAVY),
        ("SHAP", "Post Hoc\nGlobal + Local", GREEN, NAVY),
    ]
    for i, (model, typ, bg, tc) in enumerate(journey):
        x = Inches(0.4) + i * Inches(4.2)
        add_rect(slide, x, Inches(4.9), Inches(3.9), Inches(1.9), bg)
        add_text(slide, model, x + Inches(0.15), Inches(5.05),
                 Inches(3.6), Inches(0.5),
                 font_size=16, bold=True, color=tc, align=PP_ALIGN.CENTER)
        add_text(slide, typ, x + Inches(0.15), Inches(5.55),
                 Inches(3.6), Inches(0.85),
                 font_size=13, color=tc, align=PP_ALIGN.CENTER)


def slide_wolf_dog(prs):
    """Slide 7 — Wolf/Dog example"""
    slide = content_slide(prs, "High Accuracy ≠ Correct Reasoning")
    # Left: story
    add_multiline(slide, [
        "A team builds a wolf vs. dog classifier.",
        "",
        "Test accuracy:  98% ✅",
        "",
        "They deploy it. Then use an interpretation",
        "method. The discovery is surprising...",
        "",
        "The model is not looking at the animals.",
        "It learned the BACKGROUND.",
        "",
        "Wolves → snowy/forest backgrounds",
        "Dogs  → indoor/domestic backgrounds",
    ], Inches(0.4), Inches(1.3), Inches(5.8), Inches(5.7),
       font_size=15, color=DARK_GRAY)
    # Right: visual explanation
    add_rect(slide, Inches(6.5), Inches(1.3), Inches(6.5), Inches(5.7), NAVY)
    add_text(slide, "What the model ACTUALLY learned",
             Inches(6.7), Inches(1.45), Inches(6.1), Inches(0.5),
             font_size=16, bold=True, color=GREEN)
    panels = [
        ("🐺 Snow background → WOLF", "Background pixels =\nhigh activation", GREEN),
        ("🐕 Indoor background → DOG", "Animal features =\nignored", ORANGE),
    ]
    for i, (label, detail, color) in enumerate(panels):
        y = Inches(2.1) + i * Inches(2.3)
        add_rect(slide, Inches(6.7), y, Inches(6.0), Inches(2.0),
                 RGBColor(0x1E, 0x3A, 0x5F))
        add_text(slide, label, Inches(6.9), y + Inches(0.15),
                 Inches(5.6), Inches(0.5),
                 font_size=15, bold=True, color=color)
        add_text(slide, detail, Inches(6.9), y + Inches(0.75),
                 Inches(5.6), Inches(1.0),
                 font_size=13, color=WHITE)
    add_rect(slide, Inches(6.7), Inches(6.55), Inches(6.0), Inches(0.03), ORANGE)
    add_text(slide, "Without interpretability, this bug stays hidden until production fails.",
             Inches(6.7), Inches(6.6), Inches(6.0), Inches(0.45),
             font_size=13, italic=True, color=ORANGE)


def slide_why_matters(prs):
    """Slide 8 — Why interpretability matters (6 reasons)"""
    slide = content_slide(prs, "Why Interpretability Matters")
    reasons = [
        ("⚖️", "Legal Compliance",    "GDPR right to explanation (Article 22)\nfor automated decisions"),
        ("🤝", "Build Trust",         "Stakeholders accept models they\ncan understand"),
        ("🐛", "Catch Model Bugs",    "Surface unexpected correlations\nbefore deployment"),
        ("📈", "Improve Models",      "Understand failures to guide\nfeature engineering"),
        ("🎯", "Business Alignment",  "Data Scientists communicate\nfindings to non-technical teams"),
        ("⚠️", "Detect Bias",         "Identify unfair patterns\n(gender, age, race signals)"),
    ]
    for i, (icon, title, detail) in enumerate(reasons):
        col = i % 3
        row = i // 3
        x = Inches(0.35) + col * Inches(4.32)
        y = Inches(1.3)  + row * Inches(2.7)
        add_rect(slide, x, y, Inches(4.1), Inches(2.5), LIGHT_GRAY)
        add_rect(slide, x, y, Inches(4.1), Inches(0.55), NAVY)
        add_text(slide, f"{icon}  {title}",
                 x + Inches(0.15), y + Inches(0.06), Inches(3.8), Inches(0.45),
                 font_size=15, bold=True, color=GREEN)
        add_text(slide, detail,
                 x + Inches(0.15), y + Inches(0.7), Inches(3.8), Inches(1.6),
                 font_size=13, color=DARK_GRAY)


def slide_tradeoff(prs):
    """Slide 9 — Performance vs Interpretability tradeoff"""
    slide = content_slide(prs, "The Performance / Interpretability Tradeoff")
    # Left: axis diagram
    add_rect(slide, Inches(0.4), Inches(1.3), Inches(6.2), Inches(5.8), LIGHT_GRAY)

    # Axes
    add_rect(slide, Inches(1.0), Inches(6.5), Inches(5.0), Inches(0.04), DARK_GRAY)  # x-axis
    add_rect(slide, Inches(1.0), Inches(1.7), Inches(0.04), Inches(4.83), DARK_GRAY)  # y-axis
    add_text(slide, "Interpretability →",
             Inches(2.5), Inches(6.65), Inches(3.0), Inches(0.4),
             font_size=12, color=DARK_GRAY)
    add_text(slide, "Performance\n↑",
             Inches(0.45), Inches(3.5), Inches(0.5), Inches(1.0),
             font_size=11, color=DARK_GRAY)

    # Models on the curve
    models = [
        (Inches(5.5), Inches(2.2), "Neural Net\nRandom Forest\nXGBoost", ORANGE, "black box"),
        (Inches(3.5), Inches(3.2), "SVM\nk-NN", MID_GRAY, "moderate"),
        (Inches(1.8), Inches(4.5), "Decision Tree\nLinear Regression", GREEN, "interpretable"),
    ]
    for mx, my, label, color, _ in models:
        add_rect(slide, mx - Inches(0.12), my - Inches(0.12),
                 Inches(0.24), Inches(0.24), color)
        add_text(slide, label, mx - Inches(0.8), my - Inches(0.5),
                 Inches(2.0), Inches(0.9), font_size=11, color=DARK_GRAY)

    # Curve (approximated with filled rectangles)
    for i in range(20):
        t = i / 19
        cx = Inches(1.2) + t * Inches(4.6)
        cy = Inches(6.3) - t * Inches(0.5) - (1 - t) * Inches(4.0) * (1 - t)
        add_rect(slide, cx, cy, Inches(0.08), Inches(0.08), GREEN)

    add_text(slide, "SHAP bridges the gap →",
             Inches(1.3), Inches(1.8), Inches(4.5), Inches(0.4),
             font_size=12, bold=True, color=GREEN)

    # Right: the answer
    add_rect(slide, Inches(6.95), Inches(1.3), Inches(6.0), Inches(5.8), NAVY)
    add_text(slide, "Can we have both?",
             Inches(7.15), Inches(1.5), Inches(5.6), Inches(0.6),
             font_size=20, bold=True, color=WHITE)
    add_text(slide, "YES — with post hoc interpretability",
             Inches(7.15), Inches(2.2), Inches(5.6), Inches(0.55),
             font_size=18, bold=True, color=GREEN)
    add_multiline(slide, [
        "Train the most powerful model:",
        "  → XGBoost, Random Forest, etc.",
        "",
        "Apply SHAP after training:",
        "  → Get explanations without\n     sacrificing performance",
        "",
        "Industry reality (finance, insurance):",
        "  → Regulators now require explanations",
        "  → SHAP is the standard tool",
    ], Inches(7.15), Inches(2.85), Inches(5.6), Inches(4.0),
       font_size=13, color=WHITE)


def slide_technique_landscape(prs):
    """Slide 10 — Technique decision matrix"""
    slide = content_slide(prs, "Technique Landscape — When to Use What")
    headers = ["Technique", "Scope", "Model", "Output", "Best for"]
    rows = [
        ["Decision Tree viz",   "Global + Local", "Trees only",     "Tree diagram",      "Simple, shallow models"],
        ["Feature Importance",  "Global",         "Trees/Forests",  "Bar chart",         "Quick ranking (unstable)"],
        ["Coefficients",        "Global",         "Linear only",    "Bar chart",         "Linear models, direction"],
        ["PCA Biplot",          "Global",         "Any",            "Scatter plot",      "Data exploration"],
        ["SHAP beeswarm",       "Global",         "Any (Tree/NN)",  "Dot plot",          "Feature impact at scale"],
        ["SHAP force plot",     "Local",          "Any",            "Bar decomposition", "Single prediction explainer"],
        ["SHAP dependence",     "Global",         "Any",            "Scatter",           "Feature interactions"],
    ]
    col_widths = [Inches(2.5), Inches(1.7), Inches(1.7), Inches(2.0), Inches(3.5)]
    col_x = [Inches(0.35), Inches(2.9), Inches(4.65), Inches(6.4), Inches(8.45)]

    # Header row
    for j, (hdr, cw, cx) in enumerate(zip(headers, col_widths, col_x)):
        add_rect(slide, cx, Inches(1.3), cw, Inches(0.5), NAVY)
        add_text(slide, hdr, cx + Inches(0.08), Inches(1.33), cw - Inches(0.1), Inches(0.44),
                 font_size=13, bold=True, color=GREEN)

    highlight_rows = [4, 5, 6]  # SHAP rows
    for i, row in enumerate(rows):
        bg = RGBColor(0xE8, 0xFB, 0xF6) if i in highlight_rows else LIGHT_GRAY
        y = Inches(1.85) + i * Inches(0.73)
        for j, (cell, cw, cx) in enumerate(zip(row, col_widths, col_x)):
            add_rect(slide, cx, y, cw, Inches(0.7), bg)
            tc = NAVY if i in highlight_rows else DARK_GRAY
            if i in highlight_rows and j == 0:
                tc = RGBColor(0x00, 0x7A, 0x5E)
            add_text(slide, cell, cx + Inches(0.08), y + Inches(0.08),
                     cw - Inches(0.1), Inches(0.58), font_size=12, color=tc)

    add_text(slide, "★ SHAP rows highlighted — these are the focus of this masterclass",
             Inches(0.35), Inches(7.05), Inches(12.0), Inches(0.35),
             font_size=12, italic=True, color=MID_GRAY)


def slide_traditional_limits(prs):
    """Slide 11 — Traditional techniques + limitations"""
    slide = content_slide(prs, "Traditional Techniques — And Why They're Not Enough")
    techniques = [
        ("🌳 Decision Tree", [
            "Visualize yes/no decision paths",
            "Read splits from root to leaf",
            "✅ Fully transparent",
            "❌ Only shallow trees readable",
            "❌ Very unstable (tiny data change\n   → completely different tree)",
        ], Inches(0.35)),
        ("📊 Coefficients (LogReg)", [
            "Positive coef → pushes toward class 1",
            "Negative coef → pushes toward class 0",
            "✅ Direction is interpretable",
            "❌ Linear models only",
            "❌ Cannot capture interactions",
        ], Inches(4.55)),
        ("🗺️ PCA Biplot", [
            "Projects data onto 2D",
            "Arrows = feature directions",
            "✅ Good for exploration",
            "❌ Lossy projection",
            "❌ Not tied to model predictions",
        ], Inches(8.75)),
    ]
    for title, bullets, x in techniques:
        add_rect(slide, x, Inches(1.3), Inches(3.9), Inches(5.8), LIGHT_GRAY)
        add_rect(slide, x, Inches(1.3), Inches(3.9), Inches(0.55), NAVY)
        add_text(slide, title, x + Inches(0.15), Inches(1.33),
                 Inches(3.6), Inches(0.49), font_size=14, bold=True, color=GREEN)
        add_multiline(slide, bullets, x + Inches(0.2), Inches(1.95),
                      Inches(3.5), Inches(5.0), font_size=13, color=DARK_GRAY)

    add_rect(slide, Inches(0.35), Inches(7.05), Inches(12.6), Inches(0.04), ORANGE)
    add_text(slide,
             "Common weakness: all these methods fail on complex models (XGBoost, neural nets). → We need SHAP.",
             Inches(0.35), Inches(7.1), Inches(12.6), Inches(0.35),
             font_size=13, bold=True, color=ORANGE)


def slide_shap_intuition(prs):
    """Slide 12 — SHAP intuition (game theory)"""
    slide = content_slide(prs, "SHAP — The Intuition")
    # Left: game theory framing
    add_rect(slide, Inches(0.35), Inches(1.3), Inches(6.2), Inches(5.8), LIGHT_GRAY)
    add_text(slide, "Think of it as a team game",
             Inches(0.55), Inches(1.42), Inches(5.9), Inches(0.5),
             font_size=17, bold=True, color=NAVY)
    add_multiline(slide, [
        "🎮  The GAME:   one prediction for one person",
        "👥  The PLAYERS: each feature (age, income…)",
        "🏆  The SCORE:  predicted probability",
        "",
        "Question: how much did each player",
        "contribute to the final score?",
        "",
        "SHAP tries every possible combination of",
        "features and measures how much each one",
        "adds on average. That average = Shapley value.",
        "",
        "Mathematical guarantee:",
        "Σ SHAP values + expected value",
        "= actual prediction ✅",
    ], Inches(0.55), Inches(2.05), Inches(5.8), Inches(4.8),
       font_size=14, color=DARK_GRAY)

    # Right: visual decomposition
    add_rect(slide, Inches(6.85), Inches(1.3), Inches(6.1), Inches(5.8), NAVY)
    add_text(slide, "Decomposing one prediction",
             Inches(7.05), Inches(1.42), Inches(5.7), Inches(0.5),
             font_size=16, bold=True, color=WHITE)
    add_text(slide, "Baseline (avg prediction) = 0.24",
             Inches(7.05), Inches(2.05), Inches(5.7), Inches(0.45),
             font_size=14, color=MID_GRAY)

    contributions = [
        ("marital_Married", "+0.18", GREEN,  True),
        ("capital_gain",    "+0.12", GREEN,  True),
        ("age",             "+0.05", GREEN,  True),
        ("occupation",      "-0.08", ORANGE, False),
        ("gender_Female",   "-0.03", ORANGE, False),
    ]
    for i, (feat, val, color, positive) in enumerate(contributions):
        y = Inches(2.65) + i * Inches(0.65)
        bar_w = abs(float(val)) * Inches(12)
        if positive:
            bx = Inches(9.5)
        else:
            bx = Inches(9.5) - bar_w
        add_rect(slide, Inches(7.05), y, Inches(5.7), Inches(0.55),
                 RGBColor(0x1E, 0x3A, 0x5F))
        add_text(slide, feat, Inches(7.15), y + Inches(0.1),
                 Inches(2.1), Inches(0.4), font_size=12, color=WHITE)
        add_rect(slide, bx, y + Inches(0.08), bar_w + Inches(0.4), Inches(0.4), color)
        add_text(slide, val, bx + Inches(0.05), y + Inches(0.1),
                 Inches(0.6), Inches(0.38), font_size=12, bold=True, color=NAVY)

    add_rect(slide, Inches(7.05), Inches(6.0), Inches(5.7), Inches(0.04), GREEN)
    add_text(slide, "Final prediction = 0.24 + 0.18 + 0.12 + 0.05 − 0.08 − 0.03 = 0.48",
             Inches(7.05), Inches(6.1), Inches(5.7), Inches(0.45),
             font_size=13, bold=True, color=GREEN)


def slide_shap_beeswarm(prs):
    """Slide 13 — How to read SHAP beeswarm"""
    slide = content_slide(prs, "SHAP Global View — The Beeswarm Plot")
    # Left: how to read
    add_rect(slide, Inches(0.35), Inches(1.3), Inches(5.9), Inches(5.8), LIGHT_GRAY)
    add_text(slide, "📖 How to read",
             Inches(0.55), Inches(1.42), Inches(5.5), Inches(0.45),
             font_size=16, bold=True, color=NAVY)
    add_multiline(slide, [
        "Each ROW = one feature",
        "Each DOT = one person in the test set",
        "",
        "Horizontal position:",
        "  ◀ LEFT = pushes toward ≤50K (negative)",
        "  RIGHT ▶ = pushes toward >50K (positive)",
        "",
        "Dot color:",
        "  🔵 Blue = low feature value",
        "  🔴 Red  = high feature value",
        "",
        "Pattern to look for:",
        "  Red cluster RIGHT → high values push >50K",
        "  Red cluster LEFT  → high values push ≤50K",
        "",
        "Stacked dots = many people at same SHAP value",
    ], Inches(0.55), Inches(1.95), Inches(5.6), Inches(5.0),
       font_size=13, color=DARK_GRAY)

    # Right: schematic diagram
    add_rect(slide, Inches(6.55), Inches(1.3), Inches(6.4), Inches(5.8), NAVY)
    add_text(slide, "Example pattern (Census data)",
             Inches(6.75), Inches(1.42), Inches(6.0), Inches(0.45),
             font_size=14, bold=True, color=WHITE)

    features = [
        ("marital_Married",      +0.9, "High value (married) → strongly pushes >50K"),
        ("capital_gain",         +0.7, "High gain → big push, but few people have it"),
        ("educational-num",      +0.45,"More education → moderately pushes >50K"),
        ("age",                  +0.35,"Older → slight push toward >50K"),
        ("gender_Female",        -0.25,"Being female → slight push toward ≤50K ⚠️"),
        ("marital_Never-married",-0.5, "Never married → pushes toward ≤50K"),
    ]
    center_x = Inches(9.75)
    for i, (feat, shap_val, note) in enumerate(features):
        y = Inches(2.05) + i * Inches(0.77)
        # Feature label
        add_text(slide, feat, Inches(6.75), y + Inches(0.15),
                 Inches(2.4), Inches(0.4), font_size=11, color=MID_GRAY)
        # Bar
        bar_len = abs(shap_val) * Inches(1.5)
        color = GREEN if shap_val > 0 else ORANGE
        bx = center_x if shap_val > 0 else center_x - bar_len
        add_rect(slide, Inches(9.3), y + Inches(0.05), Inches(0.04), Inches(0.5), WHITE)  # center line
        add_rect(slide, bx, y + Inches(0.15), bar_len, Inches(0.3), color)
        # Note
        add_text(slide, note, Inches(11.15), y + Inches(0.1),
                 Inches(2.1), Inches(0.55), font_size=10, color=MID_GRAY, italic=True)


def slide_shap_force(prs):
    """Slide 14 — How to read SHAP force plot"""
    slide = content_slide(prs, "SHAP Local View — The Force Plot")
    # Left: explanation
    add_rect(slide, Inches(0.35), Inches(1.3), Inches(5.9), Inches(5.8), LIGHT_GRAY)
    add_text(slide, "📖 How to read",
             Inches(0.55), Inches(1.42), Inches(5.5), Inches(0.45),
             font_size=16, bold=True, color=NAVY)
    add_multiline(slide, [
        "Explains ONE specific prediction.",
        "",
        "Base value ≈ 0.24",
        "  → Model's average predicted probability",
        "  → The starting point for everyone",
        "",
        "🔴 Red blocks:",
        "  → Features pushing UP (toward >50K)",
        "  → Width = strength of push",
        "",
        "🔵 Blue blocks:",
        "  → Features pushing DOWN (toward ≤50K)",
        "  → Width = strength of pull",
        "",
        "f(x) = final predicted probability",
        "  = base + Σ(all SHAP contributions)",
    ], Inches(0.55), Inches(1.95), Inches(5.6), Inches(5.0),
       font_size=13, color=DARK_GRAY)

    # Right: visual example — Estelle
    add_rect(slide, Inches(6.55), Inches(1.3), Inches(6.4), Inches(5.8), NAVY)
    add_text(slide, "Estelle's loan application — explained",
             Inches(6.75), Inches(1.42), Inches(6.0), Inches(0.45),
             font_size=14, bold=True, color=WHITE)

    add_text(slide, "Base value: 0.24  (avg probability of >50K in population)",
             Inches(6.75), Inches(2.0), Inches(6.0), Inches(0.38),
             font_size=12, color=MID_GRAY)

    # Force bar simulation
    red_features = [("marital: Single", 0.0), ("age: 29", 0.0)]   # neutral / small
    blue_feats = [
        ("capital_gain=0",     0.45, "low savings = big negative"),
        ("education: high school", 0.25, ""),
        ("hours/wk: 38",       0.1,  ""),
    ]
    add_rect(slide, Inches(6.75), Inches(2.55), Inches(6.0), Inches(0.55), ORANGE)
    add_text(slide, "capital_gain=0  →  −0.12",
             Inches(6.9), Inches(2.62), Inches(2.8), Inches(0.4),
             font_size=13, bold=True, color=WHITE)
    add_rect(slide, Inches(9.3), Inches(2.55), Inches(3.45), Inches(0.55), RGBColor(0xFF, 0x99, 0x55))
    add_text(slide, "education  →  −0.07",
             Inches(9.4), Inches(2.62), Inches(3.2), Inches(0.4),
             font_size=13, color=WHITE)

    add_rect(slide, Inches(6.75), Inches(3.25), Inches(6.0), Inches(0.55), GREEN)
    add_text(slide, "age: 29  →  +0.02",
             Inches(6.9), Inches(3.32), Inches(5.6), Inches(0.4),
             font_size=13, bold=True, color=NAVY)

    add_rect(slide, Inches(6.75), Inches(4.0), Inches(6.0), Inches(0.04), WHITE)
    add_text(slide,
             "f(x) = 0.24 − 0.12 − 0.07 + 0.02 = 0.07",
             Inches(6.75), Inches(4.1), Inches(6.0), Inches(0.45),
             font_size=14, bold=True, color=GREEN)

    add_text(slide, "Prediction: ≤50K (7% probability of >50K)",
             Inches(6.75), Inches(4.65), Inches(6.0), Inches(0.45),
             font_size=14, color=WHITE)

    add_rect(slide, Inches(6.75), Inches(5.2), Inches(6.0), Inches(0.04), GREEN)
    add_text(slide, "Plain English explanation for Estelle:",
             Inches(6.75), Inches(5.3), Inches(6.0), Inches(0.38),
             font_size=13, bold=True, color=GREEN)
    add_text(slide,
             '"Your application was primarily affected by zero capital gain (−0.12) '
             'and education level (−0.07). Improving these factors would increase '
             'your eligibility score."',
             Inches(6.75), Inches(5.75), Inches(6.0), Inches(1.2),
             font_size=12, italic=True, color=WHITE)


def slide_shap_dependence(prs):
    """Slide 15 — SHAP dependence plot"""
    slide = content_slide(prs, "SHAP Dependence Plot — Feature Interactions")
    add_rect(slide, Inches(0.35), Inches(1.3), Inches(5.9), Inches(5.8), LIGHT_GRAY)
    add_text(slide, "📖 How to read",
             Inches(0.55), Inches(1.42), Inches(5.5), Inches(0.45),
             font_size=16, bold=True, color=NAVY)
    add_multiline(slide, [
        "X-axis: actual value of the feature",
        "Y-axis: SHAP value (contribution) for that feature",
        "Color: value of a 2nd interacting feature",
        "",
        "Upward trend → higher values push >50K",
        "Flat line → feature has no effect",
        "Color split → interaction between features",
        "",
        "Key question per plot:",
        "  At what value does SHAP cross zero?",
        "  = the tipping point for this feature",
        "",
        "Interaction signal:",
        "  If red dots cluster at one end and blue at",
        "  the other → the two features interact.",
    ], Inches(0.55), Inches(1.95), Inches(5.6), Inches(5.0),
       font_size=13, color=DARK_GRAY)

    add_rect(slide, Inches(6.55), Inches(1.3), Inches(6.4), Inches(5.8), NAVY)
    add_text(slide, "Example: age vs. SHAP value",
             Inches(6.75), Inches(1.42), Inches(6.0), Inches(0.45),
             font_size=14, bold=True, color=WHITE)

    # Schematic scatter
    import random
    random.seed(42)
    add_rect(slide, Inches(6.75), Inches(5.9), Inches(5.8), Inches(0.03), MID_GRAY)  # zero line
    add_text(slide, "SHAP=0 (baseline)", Inches(10.1), Inches(5.95),
             Inches(2.3), Inches(0.3), font_size=10, color=MID_GRAY)
    add_text(slide, "← pushes ≤50K", Inches(6.75), Inches(6.3),
             Inches(2.5), Inches(0.3), font_size=11, color=ORANGE)
    add_text(slide, "pushes >50K →", Inches(10.0), Inches(6.3),
             Inches(2.5), Inches(0.3), font_size=11, color=GREEN)

    # Dots
    for i in range(35):
        age = 18 + i * 1.8
        shap_v = -0.06 + (age - 18) * 0.0032 + random.gauss(0, 0.02)
        capital_high = random.random() > 0.6
        color = RGBColor(0xFF, 0x44, 0x44) if capital_high else RGBColor(0x44, 0x88, 0xFF)
        x_pos = Inches(6.85) + (age - 18) / 62 * Inches(5.6)
        y_pos = Inches(5.9) - shap_v * Inches(18)
        y_pos = max(Inches(2.0), min(Inches(5.85), y_pos))
        add_rect(slide, x_pos, y_pos, Inches(0.1), Inches(0.1), color)

    add_text(slide, "Age →", Inches(9.1), Inches(6.6), Inches(2.0), Inches(0.3),
             font_size=11, color=MID_GRAY)
    add_text(slide, "🔵 low capital gain   🔴 high capital gain",
             Inches(6.75), Inches(2.0), Inches(5.9), Inches(0.38),
             font_size=11, color=WHITE)
    add_text(slide,
             "Observation: SHAP for age increases with age,\nbut the effect is amplified for high capital gain (red).",
             Inches(6.75), Inches(2.5), Inches(5.9), Inches(0.8),
             font_size=12, italic=True, color=MID_GRAY)


def slide_workflow(prs):
    """Slide 16 — Practical workflow"""
    slide = content_slide(prs, "Practical Workflow — Where Interpretability Fits")
    steps = [
        ("1", "Define the problem",     "Target variable, business constraints,\nregulatory requirements",   NAVY),
        ("2", "Prepare data",           "Feature engineering, encoding,\ntrain/test split",                  NAVY),
        ("3", "Train model",            "XGBoost, Random Forest…\nOptimize for performance",                 NAVY),
        ("4", "Global SHAP",            "Beeswarm plot → top features\nCheck for bias signals",             GREEN),
        ("5", "Local SHAP",             "Force plots for key individuals\nValidate reasoning makes sense",   GREEN),
        ("6", "Communicate findings",   "Plain-English summaries for\nnon-technical stakeholders",           GREEN),
    ]
    arrow_y = Inches(4.1)
    step_w = Inches(2.0)
    gap = Inches(0.2)
    start_x = Inches(0.35)

    for i, (num, title, detail, color) in enumerate(steps):
        x = start_x + i * (step_w + gap)
        add_rect(slide, x, Inches(1.4), step_w, Inches(5.4), LIGHT_GRAY)
        add_rect(slide, x, Inches(1.4), step_w, Inches(0.55), color)
        add_text(slide, num, x + Inches(0.05), Inches(1.43),
                 Inches(0.45), Inches(0.45), font_size=22, bold=True,
                 color=GREEN if color == NAVY else NAVY, align=PP_ALIGN.CENTER)
        add_text(slide, title, x + Inches(0.1), Inches(2.05),
                 step_w - Inches(0.15), Inches(0.55),
                 font_size=13, bold=True, color=NAVY)
        add_text(slide, detail, x + Inches(0.1), Inches(2.65),
                 step_w - Inches(0.15), Inches(3.8),
                 font_size=12, color=DARK_GRAY)
        # Arrow between steps
        if i < len(steps) - 1:
            ax = x + step_w + gap * 0.1
            add_text(slide, "→", ax, arrow_y, gap * 1.5, Inches(0.45),
                     font_size=16, color=MID_GRAY, align=PP_ALIGN.CENTER)

    add_text(slide, "← Traditional techniques          SHAP (post hoc) →",
             Inches(0.35), Inches(7.05), Inches(12.5), Inches(0.38),
             font_size=13, italic=True, color=MID_GRAY)


def slide_estelle_resolved(prs):
    """Slide 17 — Estelle's case resolved with SHAP"""
    slide = content_slide(prs, "Back to Estelle — Now We Can Explain")
    # Left: force plot explanation
    add_rect(slide, Inches(0.35), Inches(1.3), Inches(7.5), Inches(5.8), LIGHT_GRAY)
    add_text(slide, "SHAP force plot for Estelle's application",
             Inches(0.55), Inches(1.42), Inches(7.1), Inches(0.45),
             font_size=15, bold=True, color=NAVY)
    add_multiline(slide, [
        "Base probability: 0.24 (population average)",
        "",
        "🔴 Factors working AGAINST loan approval:",
        "   capital_gain = 0        → −0.12",
        "   employment = 4 years    → −0.07",
        "   education: high school  → −0.05",
        "",
        "🟢 Factors SUPPORTING loan approval:",
        "   age = 29                → +0.02",
        "   hours/week = 40         → +0.01",
        "",
        "Final predicted probability: 0.03 → REJECTED",
        "",
        "→ Below the bank's threshold of 0.35",
    ], Inches(0.55), Inches(1.95), Inches(7.1), Inches(4.8),
       font_size=14, color=DARK_GRAY)

    # Right: actionable explanation
    add_rect(slide, Inches(8.15), Inches(1.3), Inches(4.8), Inches(5.8), NAVY)
    add_text(slide, "What the banker can now say:",
             Inches(8.35), Inches(1.42), Inches(4.4), Inches(0.45),
             font_size=15, bold=True, color=GREEN)
    add_text(slide,
             '"Estelle, your application was mainly '
             'affected by the absence of capital savings '
             '(no capital gain on record) and your relatively '
             'short employment history.\n\n'
             'If you can show 12+ months of savings activity '
             'and reach 5+ years of employment, '
             'your eligibility score would cross our threshold."',
             Inches(8.35), Inches(2.0), Inches(4.4), Inches(3.5),
             font_size=13, italic=True, color=WHITE)

    add_rect(slide, Inches(8.35), Inches(5.55), Inches(4.4), Inches(0.04), GREEN)
    add_text(slide, "This is why GDPR Article 22 matters.",
             Inches(8.35), Inches(5.65), Inches(4.4), Inches(0.38),
             font_size=13, color=GREEN)
    add_text(slide, "Interpretability = fairness + accountability.",
             Inches(8.35), Inches(6.1), Inches(4.4), Inches(0.4),
             font_size=13, bold=True, color=WHITE)


def slide_shap_limits(prs):
    """Slide 18 — SHAP limitations"""
    slide = content_slide(prs, "SHAP — Limitations to Know")
    limits = [
        ("🔗 Correlated features",
         "If two features are correlated (e.g. age & experience), SHAP splits credit arbitrarily. "
         "Interpretations can be misleading. Always check correlation matrix first."),
        ("⏱️ Computation cost",
         "TreeExplainer: fast (seconds). KernelExplainer (model-agnostic): can take hours on large datasets. "
         "DeepExplainer for neural nets: moderate. Choose the right explainer for your model."),
        ("⚠️ SHAP ≠ Causality",
         "A high SHAP value for 'gender' means the model uses it — not that gender causes income. "
         "Causal interpretation requires additional analysis (e.g. counterfactual methods)."),
        ("📊 Local instability",
         "Force plots for similar individuals can look very different. "
         "Don't over-interpret a single local explanation — look at patterns across many."),
    ]
    for i, (title, detail) in enumerate(limits):
        col = i % 2
        row = i // 2
        x = Inches(0.35) + col * Inches(6.45)
        y = Inches(1.3) + row * Inches(2.85)
        add_rect(slide, x, y, Inches(6.15), Inches(2.65), LIGHT_GRAY)
        add_rect(slide, x, y, Inches(6.15), Inches(0.55), ORANGE)
        add_text(slide, title, x + Inches(0.15), y + Inches(0.07),
                 Inches(5.85), Inches(0.44), font_size=15, bold=True, color=WHITE)
        add_text(slide, detail, x + Inches(0.15), y + Inches(0.65),
                 Inches(5.85), Inches(1.85), font_size=13, color=DARK_GRAY)


def slide_toolkit(prs):
    """Slide 19 — Toolkit decision table"""
    slide = content_slide(prs, "The Interpretability Toolkit — Quick Reference")
    questions = [
        ("Which features matter most (globally)?",     "SHAP beeswarm / bar chart"),
        ("How does one feature's value affect output?","SHAP dependence plot"),
        ("Why did the model predict X for person Y?",  "SHAP force plot"),
        ("What does the model look like (simple)?",    "Decision tree (if shallow)"),
        ("Quick global structure of predictions?",     "PCA biplot"),
        ("How does the model compare to a simple one?","Surrogate model (LIME/SKATER)"),
    ]
    for i, (question, answer) in enumerate(questions):
        y = Inches(1.4) + i * Inches(0.92)
        bg = RGBColor(0xE8, 0xFB, 0xF6) if i % 2 == 0 else LIGHT_GRAY
        add_rect(slide, Inches(0.35), y, Inches(12.6), Inches(0.85), bg)
        add_text(slide, f"❓  {question}",
                 Inches(0.5), y + Inches(0.15), Inches(7.5), Inches(0.6),
                 font_size=14, color=DARK_GRAY)
        add_rect(slide, Inches(8.3), y + Inches(0.1), Inches(4.3), Inches(0.65), GREEN)
        add_text(slide, f"→  {answer}",
                 Inches(8.45), y + Inches(0.15), Inches(4.0), Inches(0.55),
                 font_size=14, bold=True, color=NAVY)


def slide_conclusion(prs):
    """Slide 20 — Conclusion + next steps"""
    slide = blank_slide(prs)
    fill_bg(slide, NAVY)
    add_rect(slide, Inches(0), Inches(0), Inches(0.55), H, GREEN)
    add_rect(slide, Inches(0.7), Inches(3.6), Inches(12.3), Inches(0.04), GREEN)

    add_text(slide, "Key Takeaways",
             Inches(0.9), Inches(1.2), Inches(11.5), Inches(0.6),
             font_size=22, bold=True, color=GREEN, align=PP_ALIGN.LEFT)
    takeaways = [
        "High accuracy ≠ correct reasoning — always interpret your model",
        "SHAP is model-agnostic, theoretically grounded, and industry-standard",
        "Beeswarm for global view · Force plot for local explanations · Dependence for interactions",
        "Interpretability is not optional in regulated industries — it's a legal requirement",
        "Fairness check: always look for sensitive features in the beeswarm plot",
    ]
    for i, t in enumerate(takeaways):
        add_text(slide, f"→  {t}",
                 Inches(0.9), Inches(1.9) + i * Inches(0.38),
                 Inches(11.5), Inches(0.38),
                 font_size=14, color=WHITE)

    add_text(slide, "What's next",
             Inches(0.9), Inches(3.75), Inches(5.5), Inches(0.5),
             font_size=18, bold=True, color=GREEN)
    next_items = [
        "📘 Christoph Molnar — Interpretable ML (free): christophm.github.io/interpretable-ml-book",
        "📦 SHAP documentation: shap.readthedocs.io",
        "🔬 Try LIME for comparison: github.com/marcotcr/lime",
        "⚖️ Fairness tools: fairlearn.org · aequitas.dssg.io",
    ]
    for i, item in enumerate(next_items):
        add_text(slide, item,
                 Inches(0.9), Inches(4.35) + i * Inches(0.5),
                 Inches(11.5), Inches(0.45),
                 font_size=13, color=MID_GRAY)

    add_text(slide, "Questions & Feedback",
             Inches(0.9), Inches(6.65), Inches(6.0), Inches(0.5),
             font_size=20, bold=True, color=WHITE)
    add_text(slide, "🙋 Raise your hand in the platform",
             Inches(7.5), Inches(6.75), Inches(5.5), Inches(0.4),
             font_size=14, color=GREEN)


# ══════════════════════════════════════════════════════════════════════════════
# BUILD
# ══════════════════════════════════════════════════════════════════════════════

def build():
    prs = new_prs()

    slide_title(prs)
    slide_housekeeping(prs)
    slide_agenda(prs)

    section_divider(prs, "THE PROBLEM", "Why interpretability matters")
    slide_estelle_1(prs)
    slide_estelle_2(prs)

    section_divider(prs, "WHAT IS INTERPRETABILITY?", "Definitions & dimensions")
    slide_dimensions(prs)
    slide_wolf_dog(prs)
    slide_why_matters(prs)
    slide_tradeoff(prs)

    section_divider(prs, "TRADITIONAL TECHNIQUES", "And their limitations")
    slide_traditional_limits(prs)
    slide_technique_landscape(prs)

    section_divider(prs, "SHAP", "The modern standard")
    slide_shap_intuition(prs)
    slide_shap_beeswarm(prs)
    slide_shap_force(prs)
    slide_shap_dependence(prs)

    section_divider(prs, "PRACTICAL APPLICATION", "From model to explanation")
    slide_workflow(prs)
    slide_estelle_resolved(prs)
    slide_shap_limits(prs)
    slide_toolkit(prs)

    slide_conclusion(prs)

    out = "/Users/seb/Downloads/work_cc/slides-mc/interpretabiliy/Masterclass_Interpretability_NEW.pptx"
    prs.save(out)
    print(f"✅ Saved → {out}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    build()
